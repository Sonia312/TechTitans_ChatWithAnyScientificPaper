from flask import Flask, render_template, request, send_file
import os
from PyPDF2 import PdfReader
from werkzeug.utils import secure_filename
from docx import Document
import pptx
import textract
import pytesseract
from PIL import Image
import re
import tabula
import string
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from pathlib import Path

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])


def preprocess_data(text, image_text, tables, equations):
    # Preprocess text
    preprocessed_text = preprocess_text(text)

    # Preprocess image text
    preprocessed_image_text = preprocess_text(image_text)

    # Preprocess tables
    preprocessed_tables = preprocess_tables(tables)

    # Separate equations from text
    preprocessed_text, preprocessed_equations = separate_equations(preprocessed_text, equations)

    return preprocessed_text, preprocessed_image_text, preprocessed_tables, preprocessed_equations


def preprocess_text(text):
    # Convert text to lowercase
    text = text.lower()

    # Tokenization
    tokens = word_tokenize(text)

    # Remove punctuation
    tokens = [token for token in tokens if token not in string.punctuation]

    # Remove stopwords
    stop_words = set(stopwords.words('english'))
    tokens = [token for token in tokens if token not in stop_words]

    # Lemmatization
    lemmatizer = WordNetLemmatizer()
    tokens = [lemmatizer.lemmatize(token) for token in tokens]

    # Join tokens back into a string
    preprocessed_text = ' '.join(tokens)

    return preprocessed_text


def preprocess_tables(tables):
    # Preprocessing steps for tables (if needed)
    # Implement any specific preprocessing steps for tables here
    preprocessed_tables = tables  # No preprocessing for tables in this example
    return preprocessed_tables


def separate_equations(text, equations):
    # Separate equations from text
    preprocessed_equations = []
    for equation in equations:
        # Remove the equation from the text
        text = text.replace(equation, '')
        # Preprocess equation (if needed)
        preprocessed_equation = preprocess_equation(equation)
        preprocessed_equations.append(preprocessed_equation)
    return text, preprocessed_equations


def preprocess_equation(equation):
    # Preprocessing steps for equations (if needed)
    # Implement any specific preprocessing steps for equations here
    preprocessed_equation = equation  # No preprocessing for equations in this example
    return preprocessed_equation


def extract_text_from_pdf(pdf_path):
    with open(pdf_path, 'rb') as f:
        reader = PdfReader(f)
        text = ''
        for page_num in range(len(reader.pages)):
            text += reader.pages[page_num].extract_text()
        return text


def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = ''
    tables = []
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    # Extract tables
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)
        tables.append(table_data)
    return text, tables


def extract_text_from_pptx(pptx_path):
    ppt = pptx.Presentation(pptx_path)
    text = ''
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text


def extract_text_from_tex(tex_path):
    with open(tex_path, 'r', encoding='utf-8') as file:
        tex_content = file.read()

    # Extract equations using regular expressions
    equation_pattern = r'\$(.*?)\$'  # Assuming equations are enclosed within dollar signs ($)
    equations = re.findall(equation_pattern, tex_content)

    # Remove inline equations from the text
    text = re.sub(equation_pattern, '', tex_content)

    return text, equations


def extract_text_from_images(image_paths):
    extracted_text = ''
    for image_path in image_paths:
        try:
            image_text = pytesseract.image_to_string(Image.open(image_path))
            extracted_text += image_text + '\n'
        except OSError:
            # Ignore unsupported image files
            pass
    return extracted_text


def extract_tables_from_pdf(pdf_path):
    tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
    return tables


def extract_tables_from_docx(docx_path):
    doc = Document(docx_path)
    tables = []
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            table_data.append(row_data)
        tables.append(table_data)
    return tables


def extract_equations_from_text(text):
    # Define a regular expression pattern to match equations
    equation_pattern = r'\$.*?\$'  # Assuming equations are enclosed within dollar signs ($)

    # Find all matches of the pattern in the text
    equations = re.findall(equation_pattern, text)

    return equations


def extract_equations_from_pdf(pdf_path):
    text = extract_text_from_pdf(pdf_path)
    equations = extract_equations_from_text(text)
    return equations


def extract_equations_from_docx(docx_path):
    doc = Document(docx_path)
    equations = []
    for paragraph in doc.paragraphs:
        # Define a regular expression pattern to match equations
        equation_pattern = r'\$(.*?)\$'  # Assuming equations are enclosed within dollar signs ($)
        matches = re.findall(equation_pattern, paragraph.text)
        equations.extend(matches)
    return equations


def extract_equations_from_tex(tex_path):
    with open(tex_path, 'r', encoding='utf-8') as file:
        tex_content = file.read()

    # Extract equations using regular expressions
    equation_pattern = r'\$(.*?)\$'  # Assuming equations are enclosed within dollar signs ($)
    equations = re.findall(equation_pattern, tex_content)

    return equations


def extract_text_and_image_text_from_file(file_path):
    text = ''
    image_text = ''
    tables = []
    equations = ''
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext == '.pdf':
        text = extract_text_from_pdf(file_path)
        image_paths = [os.path.join(app.config['UPLOAD_FOLDER'], img) for img in os.listdir(app.config['UPLOAD_FOLDER'])
                       if img.endswith(('.png', '.jpg', '.jpeg', '.gif'))]
        image_text = extract_text_from_images(image_paths)
        tables = extract_tables_from_pdf(file_path)
        equations = extract_equations_from_pdf(file_path)
    elif file_ext == '.docx':
        text, tables = extract_text_from_docx(file_path)
        equations = extract_equations_from_docx(file_path)
    elif file_ext == '.pptx':
        text = extract_text_from_pptx(file_path)
    elif file_ext == '.tex':
        text, equations = extract_text_from_tex(file_path)
    elif file_ext in ['.png', '.jpg', '.jpeg', '.gif']:
        # Extract text from images using pytesseract
        image_text = extract_text_from_images([file_path])
    return text, image_text, tables, equations


@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        if 'file' not in request.files:
            return render_template('result.html', result='No file part')
        file = request.files['file']
        if file.filename == '':
            return render_template('result.html', result='No selected file')
        if file:
            upload_folder = os.path.join(os.getcwd(), app.config['UPLOAD_FOLDER'])
            if not os.path.exists(upload_folder):
                os.makedirs(upload_folder)
            filename = secure_filename(file.filename)
            filepath = os.path.join(upload_folder, filename)
            file.save(filepath)
            text, image_text, tables, equations = extract_text_and_image_text_from_file(filepath)

            # Preprocess the extracted data
            preprocessed_text, preprocessed_image_text, preprocessed_tables, preprocessed_equations = preprocess_data(
                text, image_text, tables, equations)

            # Write preprocessed data to a text file
            text_file_path = os.path.join(upload_folder, 'extracted_text.txt')
            with open(text_file_path, 'w', encoding='utf-8') as text_file:
                text_file.write(
                    f'Preprocessed Text:\n{preprocessed_text}\n\nPreprocessed Image Text:\n{preprocessed_image_text}\n\nPreprocessed Tables:\n{preprocessed_tables}\n\nPreprocessed Equations:\n{preprocessed_equations}')

            return send_file(text_file_path, as_attachment=True)
    return render_template('upload.html')


if __name__ == '__main__':
    app.run(debug=True)
